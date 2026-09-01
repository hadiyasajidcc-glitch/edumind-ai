import os
import requests
import streamlit as st
from groq import Groq
import pypdf
import docx
import pptx

# -----------------------------------------------------------------------------
# PAGE CONFIGURATION & ANIMATED PINK AESTHETIC STYLING
# -----------------------------------------------------------------------------
st.set_page_config(
    page_title="EduMind AI | Academic Hub",
    page_icon="🎓",
    layout="wide"
)

st.markdown("""
    <style>
    @keyframes fadeIn {
        from { opacity: 0; transform: translateY(12px); }
        to { opacity: 1; transform: translateY(0); }
    }

    @keyframes pulseGlow {
        0% { box-shadow: 0 0 5px rgba(255, 182, 193, 0.4); }
        50% { box-shadow: 0 0 18px rgba(199, 21, 133, 0.5); }
        100% { box-shadow: 0 0 5px rgba(255, 182, 193, 0.4); }
    }

    @keyframes floatIcon {
        0% { transform: translateY(0px); }
        50% { transform: translateY(-6px); }
        100% { transform: translateY(0px); }
    }

    .stApp {
        background: linear-gradient(135deg, #FFF0F5 0%, #FFE4E1 100%);
        color: #2D2D2D;
        animation: fadeIn 0.8s ease-in-out;
    }

    div[data-testid="stSidebar"] {
        background: linear-gradient(180deg, #FFE4E1 0%, #FFF0F5 100%);
        border-right: 2px solid #FFB6C1;
        box-shadow: 2px 0px 10px rgba(255, 182, 193, 0.3);
    }

    h1, h2, h3, .stTitle {
        color: #C71585 !important;
        font-family: 'Segoe UI', Tahoma, Geneva, Verdana, sans-serif;
        text-shadow: 1px 1px 2px rgba(255, 182, 193, 0.6);
        transition: all 0.3s ease;
    }
    
    h1:hover {
        transform: scale(1.01);
    }

    .stButton>button {
        background: linear-gradient(45deg, #D87093, #FF69B4);
        color: white;
        border-radius: 12px;
        border: none;
        padding: 0.6rem 1.2rem;
        font-weight: bold;
        transition: all 0.3s cubic-bezier(0.25, 0.8, 0.25, 1);
        box-shadow: 0 4px 10px rgba(216, 112, 147, 0.3);
    }

    .stButton>button:hover {
        background: linear-gradient(45deg, #C71585, #D87093);
        color: white;
        transform: translateY(-3px) scale(1.02);
        box-shadow: 0 8px 18px rgba(199, 21, 133, 0.4);
    }
    
    .stButton>button:active {
        transform: translateY(1px);
    }

    .stTabs [data-baseweb="tab-list"] button {
        transition: all 0.3s ease;
        border-radius: 8px;
        padding: 8px 16px;
    }
    
    .stTabs [data-baseweb="tab-list"] button:hover {
        background-color: rgba(255, 182, 193, 0.3);
        color: #C71585 !important;
    }

    .stTabs [data-baseweb="tab-list"] button[aria-selected="true"] {
        color: #C71585 !important;
        border-bottom-color: #C71585 !important;
        font-weight: bold;
        animation: pulseGlow 2s infinite;
    }

    .stTextInput>div>div>input {
        border: 2px solid #FFB6C1;
        border-radius: 10px;
        transition: all 0.3s ease;
    }

    .stTextInput>div>div>input:focus {
        border-color: #C71585;
        box-shadow: 0 0 8px rgba(199, 21, 133, 0.3);
    }

    .floating-badge {
        display: inline-block;
        animation: floatIcon 3s ease-in-out infinite;
    }
    </style>
""", unsafe_allow_html=True)

# -----------------------------------------------------------------------------
# API KEYS INITIALIZATION
# -----------------------------------------------------------------------------
GROQ_API_KEY = st.secrets.get("GROQ_API_KEY") or os.getenv("GROQ_API_KEY")
SUPADATA_API_KEY = st.secrets.get("SUPADATA_API_KEY") or os.getenv("SUPADATA_API_KEY")

groq_client = Groq(api_key=GROQ_API_KEY) if GROQ_API_KEY else None

# -----------------------------------------------------------------------------
# FAIL-PROOF INFERENCE & TRANSCRIPTION ENGINE
# -----------------------------------------------------------------------------
def get_working_models(client):
    default_models = [
        "llama-3.3-70b-versatile",
        "llama-3.1-8b-instant",
        "llama3-70b-8192",
        "llama3-8b-8192"
    ]
    try:
        live_models = client.models.list().data
        chat_models = [
            m.id for m in live_models 
            if "whisper" not in m.id.lower() and 
            any(k in m.id.lower() for k in ["llama", "mixtral", "gemma", "gpt-oss"])
        ]
        combined = chat_models + [m for m in default_models if m not in chat_models]
        return combined if combined else default_models
    except Exception:
        return default_models

def ask_groq(client, prompt):
    if not client:
        st.error("❌ Groq API key is missing. Add GROQ_API_KEY in Streamlit Secrets.")
        return "API Key Error."
        
    candidate_models = get_working_models(client)
    system_instruction = (
        "You are EduMind AI, a world-class academic tutor. "
        "Present information cleanly with headings, key points, and clear formatting."
    )
    
    last_error = ""
    for model_name in candidate_models:
        try:
            response = client.chat.completions.create(
                model=model_name,
                messages=[
                    {"role": "system", "content": system_instruction},
                    {"role": "user", "content": prompt}
                ],
                temperature=0.5
            )
            return response.choices[0].message.content
        except Exception as e:
            last_error = str(e)
            continue
            
    st.error(f"❌ Groq API Error: Unable to communicate with models. Details: {last_error}")
    return "An error occurred while generating content."

def transcribe_voice(client, audio_bytes):
    """Transcribes user voice notes using Groq's Whisper API."""
    if not client:
        st.error("❌ Groq API Key missing.")
        return ""
    try:
        transcription = client.audio.transcriptions.create(
            file=("voice_input.wav", audio_bytes),
            model="whisper-large-v3-turbo",
            response_format="text"
        )
        return transcription.strip()
    except Exception as e:
        st.error(f"❌ Could not process audio: {e}")
        return ""

# -----------------------------------------------------------------------------
# SAFE EXTRACTOR HELPERS
# -----------------------------------------------------------------------------
def get_youtube_transcript(url):
    if not SUPADATA_API_KEY:
        st.error("❌ Supadata API key missing in Streamlit Secrets.")
        return None
    endpoint = f"https://api.supadata.ai/v1/youtube/transcript?url={url}"
    headers = {"x-api-key": SUPADATA_API_KEY}
    try:
        res = requests.get(endpoint, headers=headers)
        if res.status_code == 200:
            data = res.json()
            content = data.get("content", [])
            full_text = " ".join([item.get("text", "") for item in content])
            return full_text if full_text.strip() else None
        else:
            st.error(f"Supadata Error ({res.status_code}): {res.text}")
            return None
    except Exception as e:
        st.error(f"Failed to fetch transcript: {e}")
        return None

def extract_pdf(file):
    try:
        reader = pypdf.PdfReader(file)
        text = ""
        for page in reader.pages:
            extracted = page.extract_text()
            if extracted:
                text += extracted + "\n"
        return text if text.strip() else None
    except Exception as e:
        st.error(f"Error reading PDF: {e}")
        return None

def extract_docx(file):
    try:
        doc = docx.Document(file)
        text = "\n".join([p.text for p in doc.paragraphs if p.text.strip()])
        return text if text.strip() else None
    except Exception as e:
        st.error(f"Error reading DOCX: {e}")
        return None

def extract_pptx(file):
    try:
        prs = pptx.Presentation(file)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text += shape.text + "\n"
        return text if text.strip() else None
    except Exception as e:
        st.error(f"Error reading PPTX: {e}")
        return None

# -----------------------------------------------------------------------------
# APPLICATION USER INTERFACE
# -----------------------------------------------------------------------------
st.sidebar.markdown("<h2 class='floating-badge'>🌸 Material Source</h2>", unsafe_allow_html=True)
input_mode = st.sidebar.radio("Select Input Mode:", ["Document Upload", "YouTube Video Link"])

extracted_text = ""

if input_mode == "Document Upload":
    uploaded_file = st.sidebar.file_uploader("Upload study file (.pdf, .docx, .pptx)", type=["pdf", "docx", "pptx"])
    if uploaded_file:
        file_ext = uploaded_file.name.split(".")[-1].lower()
        with st.spinner("Extracting document text..."):
            if file_ext == "pdf":
                extracted_text = extract_pdf(uploaded_file)
            elif file_ext == "docx":
                extracted_text = extract_docx(uploaded_file)
            elif file_ext == "pptx":
                extracted_text = extract_pptx(uploaded_file)
            
            if extracted_text:
                st.sidebar.success("✨ Document loaded successfully!")
            else:
                st.sidebar.error("Could not extract readable text from file.")

elif input_mode == "YouTube Video Link":
    yt_url = st.sidebar.text_input("YouTube Lecture URL:")
    if yt_url:
        with st.spinner("Fetching transcript via Supadata..."):
            extracted_text = get_youtube_transcript(yt_url)
            if extracted_text:
                st.sidebar.success("✨ Transcript loaded successfully!")

st.markdown("<h1 class='floating-badge'>🎓 EduMind AI Study Partner</h1>", unsafe_allow_html=True)
st.write("Transform lengthy lectures, videos, and documents into instant notes & practice quizzes.")

if extracted_text:
    tab1, tab2, tab3 = st.tabs(["📝 Lecture Notes", "💬 AI Tutor", "🧩 Practice Quiz"])

    with tab1:
        st.subheader("📌 Smart Lecture Notes")
        if st.button("Generate Summary Notes"):
            with st.spinner("Summarizing material..."):
                prompt = f"Provide detailed, structured study notes with bullet points and key takeaways for this content:\n\n{extracted_text[:12000]}"
                notes = ask_groq(groq_client, prompt)
                st.markdown(notes)

    with tab2:
        st.subheader("💬 Ask Your Material Anything")
        
        # Audio recorder widget for voice messages
        audio_file = st.audio_input("🎙️ Record a Voice Message / Question:")
        user_q = st.text_input("...or type your question here:")
        
        active_question = ""
        
        if audio_file is not None:
            with st.spinner("Transcribing your voice message..."):
                audio_bytes = audio_file.read()
                active_question = transcribe_voice(groq_client, audio_bytes)
                if active_question:
                    st.info(f"🗣️ **Transcribed Question:** {active_question}")
        elif user_q:
            active_question = user_q
            
        if active_question:
            with st.spinner("Thinking..."):
                prompt = f"Context:\n{extracted_text[:10000]}\n\nQuestion: {active_question}\nAnswer concisely and accurately."
                ans = ask_groq(groq_client, prompt)
                st.write(ans)

    with tab3:
        st.subheader("🧩 Practice Quiz")
        if st.button("⚡ Generate Practice Quiz"):
            with st.spinner("Creating 5-question quiz..."):
                prompt = f"Generate a 5-question multiple choice quiz based on this text, complete with answer keys and explanations:\n\n{extracted_text[:10000]}"
                quiz = ask_groq(groq_client, prompt)
                st.markdown(quiz)
else:
    st.info("👈 Please upload a document or paste a YouTube URL in the sidebar to begin.")